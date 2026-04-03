import json
import sys
from pathlib import Path

import pandas as pd

EXTRA_SITE_PACKAGES = Path("/tmp/codex_pptx")
if EXTRA_SITE_PACKAGES.exists():
    sys.path.insert(0, str(EXTRA_SITE_PACKAGES))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "output"
DOCS_DIR = BASE_DIR / "docs"

METRICS_PATH = OUTPUT_DIR / "iwpc_model_metrics.csv"
SHAP_PATH = OUTPUT_DIR / "iwpc_shap_top_features.csv"
SUMMARY_PATH = OUTPUT_DIR / "iwpc_comparison_summary.json"
PPTX_PATH = DOCS_DIR / "iwpc_model_comparison_slides.pptx"

COMPARISON_IMG = OUTPUT_DIR / "iwpc_model_comparison.png"
SCATTER_IMG = OUTPUT_DIR / "iwpc_prediction_scatter.png"
SHAP_BAR_IMG = OUTPUT_DIR / "iwpc_shap_bar.png"
SHAP_BEESWARM_IMG = OUTPUT_DIR / "iwpc_shap_beeswarm.png"

BG = RGBColor(245, 239, 227)
INK = RGBColor(29, 27, 22)
RUST = RGBColor(166, 70, 36)
OLIVE = RGBColor(77, 91, 54)
GOLD = RGBColor(181, 138, 43)
SLATE = RGBColor(63, 74, 82)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(235, 227, 212)
PALE_OLIVE = RGBColor(230, 235, 224)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.18), Inches(0.18), Inches(12.97), Inches(7.14)
    )
    border.fill.background()
    border.line.color.rgb = SOFT
    border.line.width = Pt(1.2)


def add_textbox(slide, left, top, width, height, text="", font_size=20, color=INK, bold=False,
                font_name="Georgia", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    p.alignment = align
    return box


def add_panel(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=SOFT, radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.0)
    return shape


def add_paragraph(frame, text, font_size=18, color=INK, bold=False, level=0, space_after=6):
    p = frame.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Georgia"
    return p


def add_footer(slide, label, page_number):
    add_textbox(slide, Inches(0.55), Inches(7.02), Inches(5.0), Inches(0.25), label, 10, SLATE, False, "Calibri")
    add_textbox(slide, Inches(11.95), Inches(7.02), Inches(0.6), Inches(0.25), f"{page_number:02d}", 10, SLATE, False, "Calibri", PP_ALIGN.RIGHT)


def add_title_block(slide, eyebrow, title):
    add_textbox(slide, Inches(0.55), Inches(0.45), Inches(4.2), Inches(0.35), eyebrow.upper(), 11, RUST, False, "Calibri")
    add_textbox(slide, Inches(0.55), Inches(0.78), Inches(11.7), Inches(1.0), title, 28, INK, True)


def slide_title(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    best = metrics.iloc[0]
    iwpc = metrics.loc[metrics["model"] == "IWPC Pharmacogenetic Calculator"].iloc[0]
    clinical = metrics.loc[metrics["model"] == "IWPC Clinical Formula"].iloc[0]

    add_textbox(slide, Inches(0.55), Inches(0.46), Inches(4.5), Inches(0.3), "WARFARIN DOSE MODELING STUDY", 11, RUST, False, "Calibri")
    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.82),
        Inches(11.6),
        Inches(1.1),
        "IWPC calculator benchmarking with machine-learning baselines and explainability",
        26,
        INK,
        True,
    )

    lead_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.86), Inches(7.7), Inches(1.1))
    lead = lead_box.text_frame
    lead.word_wrap = True
    lead.text = ""
    add_paragraph(
        lead,
        "Shared 20% holdout evaluation on the IWPC cohort using the official IWPC clinical and pharmacogenetic formulas alongside linear regression, random forest, XGBoost, neural network, and the tuned LightGBM repository model.",
        font_size=17,
        color=SLATE,
        space_after=0,
    )

    panel = add_panel(slide, Inches(0.55), Inches(3.0), Inches(7.2), Inches(2.35), fill_rgb=WHITE)
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(2)
    add_textbox(slide, Inches(0.82), Inches(3.25), Inches(2.0), Inches(0.3), "BEST MODEL", 11, SLATE, False, "Calibri")
    add_textbox(slide, Inches(0.82), Inches(3.55), Inches(6.3), Inches(0.8), str(best["model"]), 28, OLIVE, True)
    add_textbox(slide, Inches(0.82), Inches(4.28), Inches(6.3), Inches(0.7), f"RMSE {best['rmse']:.2f} mg/week   |   MAE {best['mae']:.2f} mg/week   |   R² {best['r2']:.3f}", 17, INK, False, "Calibri")
    add_textbox(slide, Inches(0.82), Inches(4.72), Inches(6.3), Inches(0.4), f"Within 20% accuracy: {best['within_20_pct']:.1f}%", 17, RUST, True, "Calibri")

    note = add_panel(slide, Inches(8.05), Inches(2.55), Inches(4.35), Inches(2.8), fill_rgb=PALE_OLIVE)
    note.line.color.rgb = SOFT
    t = slide.shapes.add_textbox(Inches(8.32), Inches(2.82), Inches(3.8), Inches(2.2)).text_frame
    t.word_wrap = True
    t.text = ""
    add_paragraph(t, f"RMSE gain vs IWPC PGx: {iwpc['rmse'] - best['rmse']:.02f} mg/week", 16, INK, True)
    add_paragraph(t, f"Within-20% gain vs IWPC PGx: {best['within_20_pct'] - iwpc['within_20_pct']:.1f} points", 16, INK, True)
    add_paragraph(t, f"Clinical-only IWPC: RMSE {clinical['rmse']:.2f}, MAE {clinical['mae']:.2f}", 15, SLATE)

    add_footer(slide, "Source cohort: IWPC Subject Data", 1)


def slide_leaderboard(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(slide, "Leaderboard", "Top three models are tightly grouped, but the tuned repository model lands first on RMSE")

    card_positions = [0.55, 4.35, 8.15]
    for idx, row in enumerate(metrics.head(3).itertuples(index=False)):
        left = Inches(card_positions[idx])
        add_panel(slide, left, Inches(1.65), Inches(3.1), Inches(4.7), fill_rgb=WHITE)
        add_textbox(slide, left + Inches(0.22), Inches(1.9), Inches(0.8), Inches(0.25), f"{int(row.rank):02d}", 11, RUST, False, "Calibri")
        add_textbox(slide, left + Inches(0.22), Inches(2.18), Inches(2.6), Inches(0.8), str(row.model), 21, INK, True)
        pairs = [
            ("RMSE", f"{row.rmse:.2f}"),
            ("MAE", f"{row.mae:.2f}"),
            ("R²", f"{row.r2:.3f}"),
            ("Within 20%", f"{row.within_20_pct:.1f}%"),
        ]
        y = 3.08
        for label, value in pairs:
            add_textbox(slide, left + Inches(0.22), Inches(y), Inches(1.2), Inches(0.25), label.upper(), 10, SLATE, False, "Calibri")
            add_textbox(slide, left + Inches(0.22), Inches(y + 0.2), Inches(1.6), Inches(0.42), value, 22, OLIVE, True, "Calibri")
            y += 0.78

    add_footer(slide, "Metric priority: RMSE, then MAE", 2)


def slide_visual_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(slide, "Visual Comparison", "Aggregate metrics show only a narrow edge over the IWPC pharmacogenetic calculator")

    add_panel(slide, Inches(0.55), Inches(1.55), Inches(7.85), Inches(5.0), fill_rgb=WHITE)
    slide.shapes.add_picture(str(COMPARISON_IMG), Inches(0.78), Inches(1.78), width=Inches(7.4), height=Inches(4.55))

    insight_panel = add_panel(slide, Inches(8.65), Inches(1.55), Inches(3.8), Inches(5.0), fill_rgb=PALE_OLIVE)
    frame = slide.shapes.add_textbox(Inches(8.92), Inches(1.85), Inches(3.28), Inches(4.4)).text_frame
    frame.word_wrap = True
    frame.text = ""
    for line in [
        "The tuned LightGBM model leads on RMSE and R², but the IWPC pharmacogenetic formula remains highly competitive.",
        "Linear regression almost matches the top two, which suggests the signal is still strongly structured and not purely nonlinear.",
        "Random forest underperforms on this cohort and should not be treated as a preferred baseline here.",
        "The clinical-only IWPC formula shows a clear drop, which supports keeping genotype-aware modeling in the workflow.",
    ]:
        add_paragraph(frame, line, 16, INK, False, space_after=14)

    add_footer(slide, "Images generated from repository outputs", 3)


def slide_fit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(slide, "Fit Diagnostics", "Patient-level scatter plots show similar calibration bands for the top genotype-aware approaches")
    add_panel(slide, Inches(0.55), Inches(1.55), Inches(11.9), Inches(5.35), fill_rgb=WHITE)
    slide.shapes.add_picture(str(SCATTER_IMG), Inches(0.82), Inches(1.82), width=Inches(11.35), height=Inches(4.9))
    add_footer(slide, "Shared holdout split: random_state 42", 4)


def slide_shap(prs, shap_df, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(slide, "Explainability", "SHAP confirms that the best model still follows clinically credible genetic and anthropometric drivers")

    add_panel(slide, Inches(0.55), Inches(1.55), Inches(6.0), Inches(5.1), fill_rgb=WHITE)
    slide.shapes.add_picture(str(SHAP_BAR_IMG), Inches(0.78), Inches(1.8), width=Inches(5.55), height=Inches(4.55))

    add_panel(slide, Inches(6.8), Inches(1.55), Inches(5.65), Inches(5.1), fill_rgb=PALE_OLIVE)
    text_frame = slide.shapes.add_textbox(Inches(7.05), Inches(1.82), Inches(5.1), Inches(1.35)).text_frame
    text_frame.word_wrap = True
    text_frame.text = ""
    for line in [
        "Global feature influence is led by VKORC1, then weight, age, and CYP2C9.",
        "The strongest SHAP features align with known warfarin dose biology instead of opaque proxy variables.",
        f"Best model explained here: {summary['best_model']}.",
    ]:
        add_paragraph(text_frame, line, 16, INK, False, space_after=10)

    y = 3.25
    for row in shap_df.head(8).itertuples(index=False):
        add_textbox(slide, Inches(7.08), Inches(y), Inches(3.6), Inches(0.24), str(row.feature).replace("_", " "), 15, INK, False, "Calibri")
        add_textbox(slide, Inches(10.65), Inches(y), Inches(1.1), Inches(0.24), f"{row.mean_abs_shap:.3f}", 15, OLIVE, True, "Calibri", PP_ALIGN.RIGHT)
        y += 0.37

    add_footer(slide, "Best model explainability summary", 5)


def slide_appendix(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(slide, "Appendix", "Complete ranking table and SHAP distribution view")

    rows, cols = len(metrics) + 1, len(metrics.columns)
    table = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.62), Inches(6.15), Inches(4.95)).table
    headers = list(metrics.columns)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = SOFT
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.size = Pt(11)
        para.runs[0].font.name = "Calibri"
        para.runs[0].font.color.rgb = INK

    for row_idx, row in enumerate(metrics.itertuples(index=False), start=1):
        values = [
            int(row.rank),
            str(row.model),
            f"{row.rmse:.4f}",
            f"{row.mae:.4f}",
            f"{row.r2:.4f}",
            f"{row.within_20_pct:.4f}",
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(10)
            para.runs[0].font.name = "Calibri"
            para.runs[0].font.color.rgb = INK

    add_panel(slide, Inches(6.95), Inches(1.62), Inches(5.5), Inches(4.95), fill_rgb=WHITE)
    slide.shapes.add_picture(str(SHAP_BEESWARM_IMG), Inches(7.18), Inches(1.88), width=Inches(5.05), height=Inches(4.45))
    add_footer(slide, "Generated from iwpc_model_metrics.csv and iwpc_shap_top_features.csv", 6)


def main():
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    metrics = pd.read_csv(METRICS_PATH)
    shap_df = pd.read_csv(SHAP_PATH)
    summary = json.loads(SUMMARY_PATH.read_text())

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs, metrics)
    slide_leaderboard(prs, metrics)
    slide_visual_comparison(prs)
    slide_fit(prs)
    slide_shap(prs, shap_df, summary)
    slide_appendix(prs, metrics)

    prs.save(PPTX_PATH)
    print(f"Wrote {PPTX_PATH}")


if __name__ == "__main__":
    main()
